#!/usr/bin/env python3
"""KI-Strategie 2026 Präsentation Generator.

Generiert eine 7-Folien PPTX-Präsentation für das CEO-Meeting
"Feierabendbier mit KI-Fokus".

Verwendung:
    uv run python presentations/ki-strategie-2026/generate_slides.py

Ausgabe:
    presentations/ki-strategie-2026/ki-strategie-2026.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


# =============================================================================
# Farbschema
# =============================================================================
class Colors:
    """Professionelles Farbschema für die Präsentation."""

    PRIMARY = RGBColor(30, 58, 138)  # Dunkelblau #1E3A8A
    SECONDARY = RGBColor(100, 116, 139)  # Grau #64748B
    ACCENT = RGBColor(16, 185, 129)  # Grün #10B981
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(241, 245, 249)  # #F1F5F9
    DARK = RGBColor(30, 41, 59)  # #1E293B
    ORANGE = RGBColor(249, 115, 22)  # #F97316 für Highlights
    PURPLE = RGBColor(139, 92, 246)  # #8B5CF6 für MCP


# =============================================================================
# Hilfsfunktionen
# =============================================================================
def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    """Fügt eine Titelfolie hinzu."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Hintergrund
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = Colors.PRIMARY
    background.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = Colors.WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Untertitel
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.2), Inches(12.33), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = Colors.LIGHT_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER


def add_content_slide(prs: Presentation, title: str) -> "Slide":
    """Fügt eine Inhaltsfolie mit Titel hinzu und gibt das Slide-Objekt zurück."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Titel-Balken
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.2)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = Colors.PRIMARY
    title_bar.line.fill.background()

    # Titel-Text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = Colors.WHITE

    return slide


def add_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    fill_color: RGBColor,
    border_color: RGBColor | None = None,
) -> "Shape":
    """Fügt eine Box mit optionalem Rahmen hinzu."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int = 18,
    bold: bool = False,
    color: RGBColor = None,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """Fügt eine Textbox hinzu."""
    if color is None:
        color = Colors.DARK
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.word_wrap = True
    para = frame.paragraphs[0]
    para.text = text
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = color
    para.alignment = alignment


def add_bullet_points(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    points: list[str],
    font_size: int = 16,
    color: RGBColor = None,
) -> None:
    """Fügt Aufzählungspunkte hinzu."""
    if color is None:
        color = Colors.DARK
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            para = frame.paragraphs[0]
        else:
            para = frame.add_paragraph()
        para.text = f"• {point}"
        para.font.size = Pt(font_size)
        para.font.color.rgb = color
        para.space_after = Pt(8)


def add_arrow(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    color: RGBColor,
) -> "Shape":
    """Fügt einen Pfeil hinzu."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


# =============================================================================
# Folien-Erstellung
# =============================================================================
def create_slide_1_gdpr(prs: Presentation) -> None:
    """Folie 1: DSGVO-Compliance für KI-Anwendungen."""
    slide = add_content_slide(prs, "DSGVO-konform mit KI arbeiten")

    # Zwei-Spalten-Layout
    # Linke Spalte: Sofort umsetzbar
    add_box(slide, 0.5, 1.5, 6, 5.3, Colors.LIGHT_GRAY, Colors.ACCENT)
    add_text_box(
        slide,
        0.7,
        1.7,
        5.6,
        0.5,
        "✓ SOFORT UMSETZBAR",
        20,
        bold=True,
        color=Colors.ACCENT,
    )
    add_text_box(
        slide,
        0.7,
        2.3,
        5.6,
        0.5,
        "Azure OpenAI in der EU",
        22,
        bold=True,
        color=Colors.DARK,
    )
    add_bullet_points(
        slide,
        0.7,
        3.0,
        5.6,
        3.5,
        [
            "Dediziertes Deployment in EU-Region",
            "Daten bleiben in der EU (DSGVO-konform)",
            "KEINE Nutzung für Microsoft-Training",
            "Gleiche Preise wie OpenAI API",
            "99,9% Verfügbarkeit (SLA)",
            "SOC 2 & ISO 27001 zertifiziert",
        ],
        font_size=16,
    )

    # Rechte Spalte: Strategisch
    add_box(slide, 6.8, 1.5, 6, 5.3, Colors.LIGHT_GRAY, Colors.PURPLE)
    add_text_box(
        slide, 7.0, 1.7, 5.6, 0.5, "→ STRATEGISCH", 20, bold=True, color=Colors.PURPLE
    )
    add_text_box(
        slide, 7.0, 2.3, 5.6, 0.5, "Data Flywheel", 22, bold=True, color=Colors.DARK
    )
    add_bullet_points(
        slide,
        7.0,
        3.0,
        5.6,
        3.5,
        [
            "Kontinuierliches Lernen aus eigenen Daten",
            "Spezialisierung auf Immobilien-Domäne",
            "Eigene Modelle trainieren (Fine-Tuning)",
            "Wettbewerbsvorteil durch Datenqualität",
            "Strategisch wichtig: EU-US Beziehungen",
        ],
        font_size=16,
    )


def create_slide_2_mcp_image(prs: Presentation) -> None:
    """Folie 2: Agent/MCP-Trennung mit Bild."""
    slide = add_content_slide(prs, "KI-Agent und Werkzeuge sauber trennen")

    # Erklärungstext
    add_text_box(
        slide,
        0.5,
        1.4,
        12.33,
        0.6,
        "Das MCP-Protokoll trennt das 'Gehirn' (Agent) von den 'Werkzeugen' (Tools) – wie ein Stecker-Standard",
        font_size=18,
        color=Colors.SECONDARY,
    )

    # Bild einfügen
    image_path = Path(__file__).parent / "mcp-protocol-agent.png"
    if image_path.exists():
        # Bild zentriert einfügen
        img_left = Inches(1.5)
        img_top = Inches(2.0)
        img_width = Inches(10.33)
        slide.shapes.add_picture(str(image_path), img_left, img_top, width=img_width)
    else:
        # Fallback: Platzhalter-Box wenn Bild nicht vorhanden
        add_box(slide, 1.5, 2.0, 10.33, 4.5, Colors.LIGHT_GRAY, Colors.SECONDARY)
        add_text_box(
            slide,
            1.7,
            4.0,
            10,
            1,
            "[Bild: mcp-protocol-agent.png nicht gefunden]",
            font_size=16,
            color=Colors.SECONDARY,
            alignment=PP_ALIGN.CENTER,
        )


def create_slide_3_mcp_separation(prs: Presentation) -> None:
    """Folie 3: Agent-Trennung mit MCP-Protokoll (Drei-Box-Diagramm)."""
    slide = add_content_slide(prs, "KI-Agent und Werkzeuge sauber trennen")

    # Erklärungstext
    add_text_box(
        slide,
        0.5,
        1.4,
        12.33,
        0.6,
        "Das MCP-Protokoll trennt das 'Gehirn' (Agent) von den 'Werkzeugen' (Tools) – wie ein Stecker-Standard",
        font_size=18,
        color=Colors.SECONDARY,
    )

    # Drei Boxen: Agent | MCP | Tools
    box_top = 2.3
    box_height = 3.8

    # Agent-Box (links)
    add_box(slide, 0.5, box_top, 3.5, box_height, Colors.PRIMARY)
    add_text_box(
        slide,
        0.7,
        box_top + 0.2,
        3.1,
        0.5,
        "🧠 AGENT",
        22,
        bold=True,
        color=Colors.WHITE,
    )
    add_text_box(
        slide, 0.7, box_top + 0.8, 3.1, 0.4, '"Das Gehirn"', 16, color=Colors.LIGHT_GRAY
    )
    add_bullet_points(
        slide,
        0.7,
        box_top + 1.4,
        3.1,
        2.2,
        [
            "Rolle definiert",
            "Aufgaben kennt",
            "Regeln befolgt",
            "Entscheidungen trifft",
        ],
        font_size=14,
        color=Colors.WHITE,
    )

    # Pfeil 1
    add_arrow(slide, 4.1, box_top + 1.6, 0.8, 0.5, Colors.ACCENT)

    # MCP-Box (mitte)
    add_box(slide, 5.0, box_top, 3.3, box_height, Colors.PURPLE)
    add_text_box(
        slide, 5.2, box_top + 0.2, 2.9, 0.5, "🔌 MCP", 22, bold=True, color=Colors.WHITE
    )
    add_text_box(
        slide,
        5.2,
        box_top + 0.8,
        2.9,
        0.4,
        '"Der Stecker"',
        16,
        color=Colors.LIGHT_GRAY,
    )
    add_bullet_points(
        slide,
        5.2,
        box_top + 1.4,
        2.9,
        2.2,
        [
            "Standard-Protokoll",
            "Einheitliche Schnittstelle",
            "Beliebig erweiterbar",
            "Caching (mcp-refcache)",
        ],
        font_size=14,
        color=Colors.WHITE,
    )

    # Pfeil 2
    add_arrow(slide, 8.4, box_top + 1.6, 0.8, 0.5, Colors.ACCENT)

    # Tools-Box (rechts)
    add_box(slide, 9.3, box_top, 3.5, box_height, Colors.ACCENT)
    add_text_box(
        slide,
        9.5,
        box_top + 0.2,
        3.1,
        0.5,
        "🧰 TOOLS",
        22,
        bold=True,
        color=Colors.WHITE,
    )
    add_text_box(
        slide,
        9.5,
        box_top + 0.8,
        3.1,
        0.4,
        '"Der Werkzeugkasten"',
        16,
        color=Colors.LIGHT_GRAY,
    )
    add_bullet_points(
        slide,
        9.5,
        box_top + 1.4,
        3.1,
        2.2,
        [
            "IFC-MCP (BIM)",
            "BundesMCP (APIs)",
            "Weitere hinzufügbar",
            "Unabhängig nutzbar",
        ],
        font_size=14,
        color=Colors.WHITE,
    )

    # Fußzeile mit Vorteil
    add_box(slide, 0.5, 6.3, 12.33, 0.9, Colors.LIGHT_GRAY)
    add_text_box(
        slide,
        0.7,
        6.45,
        12,
        0.6,
        "→ Vorteil: Ein Agent kann viele Werkzeugkästen nutzen. Werkzeuge können ausgetauscht werden, ohne den Agent zu ändern.",
        font_size=16,
        color=Colors.DARK,
    )


def create_slide_4_examples(prs: Presentation) -> None:
    """Folie 4: Praxisbeispiele IFC-MCP und BundesMCP."""
    slide = add_content_slide(prs, "Konkrete Anwendungen – schon heute möglich")

    # Zwei Spalten
    col_width = 6.0
    col_height = 5.3
    col_top = 1.4

    # Linke Spalte: IFC-MCP
    add_box(
        slide, 0.5, col_top, col_width, col_height, Colors.LIGHT_GRAY, Colors.PRIMARY
    )
    add_text_box(
        slide,
        0.7,
        col_top + 0.15,
        5.6,
        0.5,
        "🏢 IFC-MCP",
        24,
        bold=True,
        color=Colors.PRIMARY,
    )
    add_text_box(
        slide,
        0.7,
        col_top + 0.65,
        5.6,
        0.5,
        "KI für Gebäudedaten (BIM/CAFM)",
        font_size=15,
        color=Colors.SECONDARY,
    )

    add_text_box(
        slide, 0.7, col_top + 1.15, 5.6, 0.4, "Der Agent kann fragen:", 13, bold=True
    )
    add_bullet_points(
        slide,
        0.7,
        col_top + 1.5,
        5.6,
        1.3,
        [
            '"Wie viele Türen hat das Gebäude?"',
            '"Berechne Flächen nach DIN 277"',
            '"Kostengruppen nach DIN 276?"',
        ],
        font_size=13,
    )

    add_text_box(
        slide,
        0.7,
        col_top + 2.9,
        5.6,
        0.4,
        "Features:",
        13,
        bold=True,
        color=Colors.ACCENT,
    )
    add_bullet_points(
        slide,
        0.7,
        col_top + 3.25,
        5.6,
        2.0,
        [
            "30+ spezialisierte BIM-Tools",
            "DIN 276 Kostengruppen",
            "DIN 277 Flächenberechnung",
            "VDI-konforme Energiesimulation (bim2sim)",
            "CAFM-Integration",
        ],
        font_size=12,
        color=Colors.DARK,
    )

    # Rechte Spalte: BundesMCP
    add_box(
        slide, 6.8, col_top, col_width, col_height, Colors.LIGHT_GRAY, Colors.ACCENT
    )
    add_text_box(
        slide,
        7.0,
        col_top + 0.15,
        5.6,
        0.5,
        "🗺️ BundesMCP",
        24,
        bold=True,
        color=Colors.ACCENT,
    )
    add_text_box(
        slide,
        7.0,
        col_top + 0.65,
        5.6,
        0.5,
        "Agent-Wrapper für alle deutschen Behörden-APIs",
        font_size=15,
        color=Colors.SECONDARY,
    )

    add_text_box(
        slide, 7.0, col_top + 1.15, 5.6, 0.4, "Der Agent kann fragen:", 13, bold=True
    )
    add_bullet_points(
        slide,
        7.0,
        col_top + 1.5,
        5.6,
        1.3,
        [
            '"Gibt es Hochwasserrisiko hier?"',
            '"Wie ist die Luftqualität?"',
            '"Welche Schulen sind in der Nähe?"',
        ],
        font_size=13,
    )

    add_text_box(
        slide,
        7.0,
        col_top + 2.9,
        5.6,
        0.4,
        "Features:",
        13,
        bold=True,
        color=Colors.PRIMARY,
    )
    add_bullet_points(
        slide,
        7.0,
        col_top + 3.25,
        5.6,
        2.0,
        [
            "60+ bundesAPI-Dienste",
            "OpenStreetMap-Integration",
            "Hochwasser, Luftqualität, Demographie",
            "Geocoding & Routing",
            "Automatische Standortbewertung",
        ],
        font_size=12,
        color=Colors.DARK,
    )

    # Footer
    add_text_box(
        slide,
        0.5,
        6.85,
        12.33,
        0.4,
        "→ Beide Tools nutzen mcp-refcache für effizientes Caching großer Datenmengen",
        font_size=13,
        color=Colors.SECONDARY,
        alignment=PP_ALIGN.CENTER,
    )


def create_slide_5_flowise(prs: Presentation) -> None:
    """Folie 5: Flowise AI für Workflow-Orchestrierung."""
    slide = add_content_slide(prs, "Flowise AI – KI-Workflows ohne Programmieren")

    # Hauptbereich mit drei Vorteilen
    add_text_box(
        slide,
        0.5,
        1.4,
        12.33,
        0.6,
        "Visual Builder für KI-Anwendungen: Drag & Drop statt Code",
        font_size=18,
        color=Colors.SECONDARY,
    )

    # Drei Feature-Boxen
    box_width = 3.9
    box_height = 3.2
    box_top = 2.2

    # Box 1: Visuelle Erstellung
    add_box(
        slide, 0.5, box_top, box_width, box_height, Colors.LIGHT_GRAY, Colors.PRIMARY
    )
    add_text_box(
        slide,
        0.7,
        box_top + 0.2,
        3.5,
        0.5,
        "🎨 Visuell",
        20,
        bold=True,
        color=Colors.PRIMARY,
    )
    add_bullet_points(
        slide,
        0.7,
        box_top + 0.9,
        3.5,
        2.2,
        [
            "Drag & Drop Bausteine",
            "Kein Code nötig",
            "Sofort sichtbar",
            "Schnelles Prototyping",
        ],
        font_size=14,
    )

    # Box 2: Multi-Agent
    add_box(
        slide, 4.7, box_top, box_width, box_height, Colors.LIGHT_GRAY, Colors.PURPLE
    )
    add_text_box(
        slide,
        4.9,
        box_top + 0.2,
        3.5,
        0.5,
        "🤝 Multi-Agent",
        20,
        bold=True,
        color=Colors.PURPLE,
    )
    add_bullet_points(
        slide,
        4.9,
        box_top + 0.9,
        3.5,
        2.2,
        [
            "Mehrere KI-Agenten",
            "Spezialisierte Rollen",
            "Zusammenarbeit",
            "Komplexe Aufgaben",
        ],
        font_size=14,
    )

    # Box 3: Integration
    add_box(
        slide, 8.9, box_top, box_width, box_height, Colors.LIGHT_GRAY, Colors.ACCENT
    )
    add_text_box(
        slide,
        9.1,
        box_top + 0.2,
        3.5,
        0.5,
        "🔗 Integration",
        20,
        bold=True,
        color=Colors.ACCENT,
    )
    add_bullet_points(
        slide,
        9.1,
        box_top + 0.9,
        3.5,
        2.2,
        [
            "APIs anbinden",
            "Datenbanken",
            "Dokumente",
            "Eigene Tools (MCP)",
        ],
        font_size=14,
    )

    # Vorschlag-Box
    add_box(slide, 0.5, 5.7, 12.33, 1.3, Colors.PRIMARY)
    add_text_box(
        slide,
        0.7,
        5.85,
        11.9,
        0.5,
        "💡 Vorschlag: Flowise als zentrale Plattform für KI-Workflow-Entwicklung",
        font_size=18,
        bold=True,
        color=Colors.WHITE,
    )
    add_text_box(
        slide,
        0.7,
        6.35,
        11.9,
        0.5,
        "Open Source, Self-Hosted möglich, Active Community, TypeScript-basiert",
        font_size=14,
        color=Colors.LIGHT_GRAY,
    )


# =============================================================================
# Hauptfunktion
# =============================================================================
def create_presentation() -> None:
    """Erstellt die vollständige Präsentation."""
    prs = Presentation()

    # 16:9 Format
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Titelfolie
    add_title_slide(
        prs,
        "KI-Strategie 2026",
        "Feierabendbier mit KI-Fokus • DSGVO • MCP • Flowise • Praxisbeispiele",
    )

    # Folie 1: MCP mit Bild
    create_slide_2_mcp_image(prs)

    # Folie 2: MCP Drei-Box-Diagramm
    create_slide_3_mcp_separation(prs)

    # Folie 3: Praxisbeispiele
    create_slide_4_examples(prs)

    # Folie 4: Flowise AI
    create_slide_5_flowise(prs)

    # Folie 5: DSGVO-Compliance (jetzt letzte Folie)
    create_slide_1_gdpr(prs)

    # Speichern
    output_path = Path(__file__).parent / "ki-strategie-2026.pptx"
    prs.save(str(output_path))
    print(f"✅ Präsentation erstellt: {output_path}")
    print(f"   {len(prs.slides)} Folien generiert")
    print("\nFolien-Reihenfolge:")
    print("   0. Titelfolie: KI-Strategie 2026")
    print("   1. KI-Agent und Werkzeuge (MCP-Bild)")
    print("   2. KI-Agent und Werkzeuge (Drei-Box-Diagramm)")
    print("   3. Praxisbeispiele: IFC-MCP & BundesMCP")
    print("   4. Flowise AI – KI-Workflows")
    print("   5. DSGVO-konform mit KI arbeiten")


if __name__ == "__main__":
    create_presentation()
